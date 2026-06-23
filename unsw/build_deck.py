#!/usr/bin/env python3
"""
Rare Diseases NSW Co-Design Workshop — facilitator deck builder.
Builds a 16:9 PowerPoint using python-pptx with native vector diagrams.
Framework content: Emma Blomkamp / New Know How systemic co-design.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Palette (echoes the Practice Wheel's category colours)
# ----------------------------------------------------------------------------
INK       = RGBColor(0x1A, 0x22, 0x30)   # primary text / dark anchor
INK2      = RGBColor(0x49, 0x54, 0x68)   # secondary text
PAPER     = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_ALT = RGBColor(0xF5, 0xF6, 0xF8)
LINE      = RGBColor(0xE2, 0xE5, 0xEA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

NAVY      = RGBColor(0x15, 0x24, 0x3A)   # deep anchor for title / dividers
NAVY2     = RGBColor(0x21, 0x35, 0x52)

PEOPLE    = RGBColor(0xE0, 0xA8, 0x2E)   # gold
PROCESS   = RGBColor(0x1F, 0xA6, 0x4A)   # green
PRACTICE  = RGBColor(0xE8, 0x59, 0x1C)   # orange
PLACE     = RGBColor(0x7A, 0x45, 0xD6)   # purple

# soft tints for card backgrounds
PEOPLE_T   = RGBColor(0xFB, 0xF3, 0xDD)
PROCESS_T  = RGBColor(0xE5, 0xF5, 0xEA)
PRACTICE_T = RGBColor(0xFC, 0xEC, 0xE2)
PLACE_T    = RGBColor(0xEF, 0xE8, 0xFB)

# maturity model
EXPLORE    = RGBColor(0x4A, 0x14, 0x14)
UNDERSTAND = RGBColor(0xE8, 0x59, 0x1C)
APPLY      = RGBColor(0x7A, 0x45, 0xD6)
INTEGRATE  = RGBColor(0x1F, 0xA6, 0x4A)
FLOURISH   = RGBColor(0xE0, 0xA8, 0x2E)

FONT_H = "Lexend"        # heading (accessible; substitutes cleanly if absent)
FONT_B = "Source Sans 3" # body

EMU_IN = 914400
SW, SH = 13.333, 7.5
MX = 0.72  # left/right margin

# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def _noline(shape):
    shape.line.fill.background()


def _noshadow(shape):
    shape.shadow.inherit = False


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    _noshadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        _noline(shp)
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def oval(s, x, y, w, h, fill=None, line=None, line_w=1.0, alpha=None):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    _noshadow(shp)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(shp.fill.fore_color, alpha)
    if line is None:
        _noline(shp)
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def _alpha(fore_color, opacity):
    """opacity 0..1"""
    srgb = fore_color._xFill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity * 100000))})
    srgb.append(a)


def line_h(s, x, y, w, color=LINE, weight=1.0):
    cn = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    cn.line.color.rgb = color; cn.line.width = Pt(weight)
    _noshadow(cn)
    return cn


def tb(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {t, size, bold, color, font, align, sa, sb, line,
       bullet, bcolor, italic}. 't' may be a string or list of (text, optsdict) runs."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if p.get('sa') is not None: para.space_after = Pt(p['sa'])
        if p.get('sb') is not None: para.space_before = Pt(p['sb'])
        if p.get('line') is not None: para.line_spacing = p['line']
        content = p['t']
        bullet = p.get('bullet')
        if bullet:
            r = para.add_run(); r.text = "•  "
            r.font.size = Pt(p.get('size', 16)); r.font.bold = True
            r.font.color.rgb = p.get('bcolor', INK); r.font.name = p.get('font', FONT_B)
        if isinstance(content, str):
            content = [(content, {})]
        for text, opt in content:
            r = para.add_run(); r.text = text
            r.font.size = Pt(opt.get('size', p.get('size', 16)))
            r.font.bold = opt.get('bold', p.get('bold', False))
            r.font.italic = opt.get('italic', p.get('italic', False))
            r.font.name = opt.get('font', p.get('font', FONT_B))
            r.font.color.rgb = opt.get('color', p.get('color', INK))
    return box


def kicker_title(s, kicker, title, accent=NAVY, title_color=INK, top=0.62):
    """Standard content-slide header with a coloured accent tab + kicker + title."""
    rect(s, MX, top, 0.34, 0.34, fill=accent, rounded=True, radius=0.28)
    tb(s, MX + 0.5, top - 0.06, 9.5, 0.4,
       [{'t': kicker.upper(), 'size': 13.5, 'bold': True, 'color': accent,
         'font': FONT_H}])
    tb(s, MX, top + 0.42, SW - 2 * MX, 1.0,
       [{'t': title, 'size': 32, 'bold': True, 'color': title_color, 'font': FONT_H,
         'line': 1.02}])


def page_no(s, n):
    tb(s, SW - 1.4, SH - 0.5, 0.9, 0.3,
       [{'t': str(n), 'size': 11, 'color': INK2, 'align': PP_ALIGN.RIGHT,
         'font': FONT_B}])


def footer_brand(s, color=INK2):
    tb(s, MX, SH - 0.5, 8.0, 0.3,
       [{'t': "Rare Diseases NSW Co-Design Workshop", 'size': 10.5, 'color': color,
         'font': FONT_B}])


def source_note(s, txt, color=INK2):
    tb(s, MX, SH - 0.74, 8.5, 0.3,
       [{'t': txt, 'size': 10, 'italic': True, 'color': color, 'font': FONT_B}])


N = 0
def num():
    global N
    N += 1
    return N

# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = slide(NAVY)
# colour spectrum bar (the four domains)
bx, bw = 0, SW / 4
for c in (PEOPLE, PROCESS, PRACTICE, PLACE):
    rect(s, bx, 0, bw, 0.16, fill=c); bx += bw
tb(s, MX, 1.55, 11.0, 0.5,
   [{'t': "RARE NSW  ·  CO-DESIGN WORKSHOP", 'size': 15, 'bold': True,
     'color': PEOPLE, 'font': FONT_H}])
tb(s, MX, 2.15, 11.6, 2.6,
   [{'t': "Designing Better,", 'size': 58, 'bold': True, 'color': WHITE,
     'font': FONT_H, 'line': 1.0, 'sa': 0},
    {'t': "Together", 'size': 58, 'bold': True, 'color': PEOPLE,
     'font': FONT_H, 'line': 1.0}])
tb(s, MX, 4.5, 10.8, 0.8,
   [{'t': "A co-design workshop for the rare disease community in New South Wales",
     'size': 19, 'color': RGBColor(0xCE, 0xD6, 0xE2), 'font': FONT_B, 'line': 1.2}])
line_h(s, MX, 5.7, 4.2, color=NAVY2, weight=1.5)
tb(s, MX, 5.9, 11.0, 0.9,
   [{'t': [("Facilitated by Rare NSW · Kris Pierce Consulting", {'bold': True, 'color': WHITE})],
     'size': 14, 'font': FONT_B, 'sa': 2},
    {'t': "Framework: The Systemic Design Practice — Emma Blomkamp / New Know How",
     'size': 12.5, 'color': RGBColor(0x9D, 0xAA, 0xBD), 'font': FONT_B}])

# ============================================================================
# SLIDE 2 — ACKNOWLEDGEMENT OF COUNTRY
# ============================================================================
s = slide(PAPER_ALT)
rect(s, 0, 0, 0.28, SH, fill=PROCESS)
tb(s, MX + 0.4, 1.4, 11.0, 0.4,
   [{'t': "ACKNOWLEDGEMENT OF COUNTRY", 'size': 14, 'bold': True, 'color': PROCESS,
     'font': FONT_H}])
tb(s, MX + 0.4, 2.05, 11.2, 3.2,
   [{'t': "We acknowledge the Traditional Custodians of the lands on which we meet, "
          "and pay our respects to Elders past and present.",
     'size': 27, 'bold': True, 'color': INK, 'font': FONT_H, 'line': 1.18, 'sa': 16},
    {'t': "We recognise that caring for community, listening deeply, and making "
          "decisions together are practices First Nations peoples have held for "
          "tens of thousands of years — and that good co-design has much to learn "
          "from them.",
     'size': 16, 'color': INK2, 'font': FONT_B, 'line': 1.3}])
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 3 — AIM OF THE DAY
# ============================================================================
s = slide()
kicker_title(s, "Welcome", "Why we're here today", accent=PROCESS)
tb(s, MX, 1.95, 6.0, 0.9,
   [{'t': "Rare conditions are individually rare but collectively common — and the "
          "people who live with them hold knowledge no system can design without.",
     'size': 16.5, 'color': INK2, 'font': FONT_B, 'line': 1.3}])
tb(s, MX, 3.05, 6.0, 0.5,
   [{'t': "Our aim for the day", 'size': 16, 'bold': True, 'color': INK, 'font': FONT_H}])
aims = [
    "Build a shared language for how we work together — co-design.",
    "Pressure-test an established framework against the reality of rare disease.",
    "Agree how we'll behave together — our Code of Care.",
    "Leave with first, concrete steps for Rare Diseases NSW.",
]
yy = 3.55
for a in aims:
    oval(s, MX, yy + 0.05, 0.16, 0.16, fill=PROCESS)
    tb(s, MX + 0.34, yy - 0.04, 5.7, 0.6,
       [{'t': a, 'size': 14.5, 'color': INK, 'font': FONT_B, 'line': 1.18}])
    yy += 0.74
# right "success looks like" panel
px = 7.4
rect(s, px, 1.95, SW - MX - px, 4.5, fill=NAVY, rounded=True, radius=0.04)
tb(s, px + 0.45, 2.35, SW - MX - px - 0.9, 0.5,
   [{'t': "SUCCESS LOOKS LIKE", 'size': 13, 'bold': True, 'color': PEOPLE, 'font': FONT_H}])
succ = [
    ("Everyone contributes", "Lived experience, practice wisdom and evidence all in the room."),
    ("Honest, not polished", "We name what doesn't fit for rare disease."),
    ("A framework that's ours", "Adapted, not adopted off the shelf."),
    ("Momentum", "Clear owners and next steps before we leave."),
]
yy = 2.95
for t, d in succ:
    tb(s, px + 0.45, yy, SW - MX - px - 0.9, 0.8,
       [{'t': [("→  ", {'color': PEOPLE, 'bold': True}), (t, {'bold': True, 'color': WHITE})],
         'size': 14.5, 'font': FONT_H, 'sa': 1},
        {'t': d, 'size': 12.5, 'color': RGBColor(0xBE, 0xC8, 0xD6), 'font': FONT_B, 'line': 1.15}])
    yy += 0.86
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 4 — HOW TODAY FLOWS (AGENDA)
# ============================================================================
s = slide()
kicker_title(s, "Agenda", "How today flows", accent=NAVY)
steps = [
    ("01", "Ground", "Welcome, aim of the day, and what we mean by co-design.", PEOPLE),
    ("02", "Present", "Walk through the current framework — principles & the four domains.", PROCESS),
    ("03", "Test", "Do these make sense? What's missing for rare disease? — discussion.", PRACTICE),
    ("04", "Agree", "Co-create our Code of Care and map where we are.", PLACE),
    ("05", "Forward", "First steps for Rare Diseases NSW — the 'how to'.", NAVY),
]
colw = (SW - 2 * MX - 4 * 0.2) / 5
x = MX
for n, t, d, c in steps:
    top = 2.3
    rect(s, x, top, colw, 3.4, fill=PAPER_ALT, rounded=True, radius=0.05)
    rect(s, x, top, colw, 0.14, fill=c, rounded=False)
    tb(s, x + 0.22, top + 0.32, colw - 0.44, 0.7,
       [{'t': n, 'size': 30, 'bold': True, 'color': c, 'font': FONT_H}])
    tb(s, x + 0.22, top + 1.05, colw - 0.44, 0.5,
       [{'t': t, 'size': 16.5, 'bold': True, 'color': INK, 'font': FONT_H}])
    line_h(s, x + 0.22, top + 1.55, colw - 0.44, color=LINE, weight=1.0)
    tb(s, x + 0.22, top + 1.7, colw - 0.44, 1.6,
       [{'t': d, 'size': 12.5, 'color': INK2, 'font': FONT_B, 'line': 1.22}])
    x += colw + 0.2
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 5 — OUR SUGGESTED APPROACH
# ============================================================================
s = slide(NAVY)
rect(s, 0, 0, SW, 0.16, fill=PROCESS)
tb(s, MX, 0.7, 11.0, 0.4,
   [{'t': "OUR SUGGESTED APPROACH", 'size': 13.5, 'bold': True, 'color': PEOPLE, 'font': FONT_H}])
tb(s, MX, 1.12, 11.6, 0.9,
   [{'t': "Start from what exists — then make it fit rare disease",
     'size': 30, 'bold': True, 'color': WHITE, 'font': FONT_H, 'line': 1.04}])
appr = [
    ("Present", "Share the current co-design framework — its principles, values and four domains.", PEOPLE),
    ("Test", "Hold it up against rare disease reality. What resonates? What's missing?", PROCESS),
    ("Co-create", "Adapt the framework and agree how we'll work together — our Code of Care.", PRACTICE),
    ("Map & move", "Locate where we are on the journey, and commit to first steps.", PLACE),
]
cw = (SW - 2 * MX - 3 * 0.28) / 4
x = MX
for i, (t, d, c) in enumerate(appr):
    top = 2.65
    rect(s, x, top, cw, 3.1, fill=NAVY2, rounded=True, radius=0.05)
    rect(s, x, top, 0.5, 0.5, fill=c, rounded=True, radius=0.2)
    tb(s, x, top, 0.5, 0.5, [{'t': str(i + 1), 'size': 18, 'bold': True, 'color': WHITE,
                              'align': PP_ALIGN.CENTER, 'font': FONT_H}],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(s, x + 0.28, top + 0.72, cw - 0.5, 0.5,
       [{'t': t, 'size': 19, 'bold': True, 'color': WHITE, 'font': FONT_H}])
    tb(s, x + 0.28, top + 1.3, cw - 0.56, 1.7,
       [{'t': d, 'size': 13, 'color': RGBColor(0xC2, 0xCC, 0xDA), 'font': FONT_B, 'line': 1.25}])
    if i < 3:
        tb(s, x + cw - 0.02, top + 1.2, 0.3, 0.5,
           [{'t': "›", 'size': 26, 'bold': True, 'color': PEOPLE, 'align': PP_ALIGN.CENTER,
             'font': FONT_H}])
    x += cw + 0.28
tb(s, MX, 6.25, 11.8, 0.5,
   [{'t': [("Throughout: ", {'bold': True, 'color': PEOPLE}),
           ("nothing about us, without us. People with lived experience lead, not just inform.",
            {'color': RGBColor(0xC2, 0xCC, 0xDA)})],
     'size': 13.5, 'font': FONT_B, 'italic': True}])
page_no(s, num())

# ============================================================================
# SLIDE 6 — WHAT WE MEAN BY CO-DESIGN
# ============================================================================
s = slide()
kicker_title(s, "Shared language", "What we mean by co-design", accent=PEOPLE, title_color=INK)
tb(s, MX, 2.0, SW - 2 * MX, 1.0,
   [{'t': [("Co-design is ", {}),
           ("designing with people, not for them", {'bold': True, 'color': PRACTICE}),
           (" — sharing power and decisions with those who have lived experience of "
            "the issue, alongside practice wisdom, evidence and cultural knowledge.", {})],
     'size': 22, 'color': INK, 'font': FONT_H, 'line': 1.25}])
pairs = [
    ("It is", [
        "Sharing power and decisions",
        "Many kinds of knowledge, valued equally",
        "A mindset and a set of practices",
        "Slower at the start, stronger in the end",
    ], PROCESS, PROCESS_T),
    ("It is not", [
        "Consultation or a feedback survey",
        "Experts deciding, then 'checking in'",
        "A single workshop or a tick-box",
        "Co-design in name only ('co-washing')",
    ], PRACTICE, PRACTICE_T),
]
x = MX
cw = (SW - 2 * MX - 0.5) / 2
for title, items, c, t in pairs:
    top = 3.4
    rect(s, x, top, cw, 3.0, fill=t, rounded=True, radius=0.05)
    rect(s, x, top, 0.16, 3.0, fill=c)
    tb(s, x + 0.45, top + 0.28, cw - 0.7, 0.5,
       [{'t': title.upper(), 'size': 15, 'bold': True, 'color': c, 'font': FONT_H}])
    yy = top + 0.95
    for it in items:
        tb(s, x + 0.45, yy, cw - 0.9, 0.5,
           [{'t': [("•  ", {'color': c, 'bold': True}), (it, {'color': INK})],
             'size': 14, 'font': FONT_B, 'line': 1.1}])
        yy += 0.5
    x += cw + 0.5
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 7 — THE SYSTEMIC DESIGN PRACTICE WHEEL (overview)
# ============================================================================
s = slide()
kicker_title(s, "The current framework", "The Systemic Design Practice", accent=NAVY)
tb(s, MX, 1.95, 5.4, 1.4,
   [{'t': "One practical map for navigating complex, systemic challenges. Four "
          "domains to work through, held together by six guiding principles.",
     'size': 15, 'color': INK2, 'font': FONT_B, 'line': 1.28}])
# six principles list (left)
tb(s, MX, 3.35, 5.4, 0.4,
   [{'t': "SIX GUIDING PRINCIPLES", 'size': 12.5, 'bold': True, 'color': INK, 'font': FONT_H}])
principles = ["Purpose driven", "Inclusive collaboration", "Recognising complexity",
              "Adaptive learning", "Equalising power", "Self-determination"]
yy = 3.78
pcols = [PROCESS, PEOPLE, PRACTICE, PROCESS, PLACE, PRACTICE]
for i, p in enumerate(principles):
    oval(s, MX, yy + 0.04, 0.13, 0.13, fill=pcols[i])
    tb(s, MX + 0.28, yy - 0.05, 5.0, 0.4,
       [{'t': p, 'size': 14, 'bold': True, 'color': INK, 'font': FONT_B}])
    yy += 0.42
# 2x2 quadrant grid (right)
gx, gy = 6.55, 1.95
gw, gh = SW - MX - gx, 4.5
cellw = (gw - 0.18) / 2
cellh = (gh - 0.18) / 2
quad = [
    ("PEOPLE", "Who needs to be involved?", ["Self", "Team", "Organisation", "Community"], PEOPLE, PEOPLE_T, INK),
    ("PROCESS", "How will we approach this?", ["Invitation", "Preparation", "Navigation", "Completion"], PROCESS, PROCESS_T, INK),
    ("PLACE", "Where does this fit?", ["Symbols", "Products", "Interactions", "System"], PLACE, PLACE_T, INK),
    ("PRACTICE", "What will we use?", ["Tools", "Techniques", "Methods", "Principles"], PRACTICE, PRACTICE_T, INK),
]
pos = [(gx, gy), (gx + cellw + 0.18, gy), (gx, gy + cellh + 0.18), (gx + cellw + 0.18, gy + cellh + 0.18)]
for (name, q, items, c, t, tc), (cx, cy) in zip(quad, pos):
    rect(s, cx, cy, cellw, cellh, fill=t, rounded=True, radius=0.06)
    rect(s, cx, cy, cellw, 0.12, fill=c)
    tb(s, cx + 0.24, cy + 0.22, cellw - 0.4, 0.4,
       [{'t': name, 'size': 15, 'bold': True, 'color': c, 'font': FONT_H}])
    tb(s, cx + 0.24, cy + 0.6, cellw - 0.4, 0.4,
       [{'t': q, 'size': 11, 'italic': True, 'color': INK2, 'font': FONT_B}])
    tb(s, cx + 0.24, cy + 1.02, cellw - 0.4, 1.0,
       [{'t': "  ·  ".join(items), 'size': 11.5, 'color': INK, 'font': FONT_B, 'line': 1.2}])
source_note(s, "Framework: Emma Blomkamp, The Systemic Design Practice Wheel (2021), CC BY-NC.")
page_no(s, num())

# ============================================================================
# SLIDE 8 — PRINCIPLES & VALUES
# ============================================================================
s = slide()
kicker_title(s, "The current framework", "Principles & values", accent=PROCESS)
cards = [
    ("Purpose driven", "Stay anchored to why this matters and who it's for.", PROCESS),
    ("Inclusive collaboration", "Bring diverse people and knowledge together — genuinely.", PEOPLE),
    ("Recognising complexity", "Work with the whole system, not isolated parts.", PRACTICE),
    ("Adaptive learning", "Stay curious; test, learn and adjust as we go.", PROCESS),
    ("Equalising power", "Share decisions; redistribute who gets to shape outcomes.", PLACE),
    ("Self-determination", "People shape the things that affect their own lives.", PRACTICE),
]
cw = (SW - 2 * MX - 2 * 0.3) / 3
ch = 1.78
positions = []
for r in range(2):
    for col in range(3):
        positions.append((MX + col * (cw + 0.3), 2.05 + r * (ch + 0.3)))
for (t, d, c), (x, y) in zip(cards, positions):
    rect(s, x, y, cw, ch, fill=PAPER_ALT, rounded=True, radius=0.06)
    rect(s, x, y, 0.14, ch, fill=c)
    tb(s, x + 0.36, y + 0.26, cw - 0.55, 0.5,
       [{'t': t, 'size': 16.5, 'bold': True, 'color': INK, 'font': FONT_H}])
    tb(s, x + 0.36, y + 0.78, cw - 0.6, 0.9,
       [{'t': d, 'size': 13, 'color': INK2, 'font': FONT_B, 'line': 1.22}])
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 9 — DISCUSSION: do these make sense? what's missing?
# ============================================================================
s = slide(NAVY)
rect(s, 0, 0, SW, 0.16, fill=PEOPLE)
tb(s, MX, 0.62, 11.0, 0.4,
   [{'t': "DISCUSSION", 'size': 13.5, 'bold': True, 'color': PEOPLE, 'font': FONT_H}])
tb(s, MX, 1.02, 11.8, 0.9,
   [{'t': "Do these make sense for rare disease? What's missing?",
     'size': 29, 'bold': True, 'color': WHITE, 'font': FONT_H, 'line': 1.05}])
# left: the two questions
qx = MX
cw = 5.55
q1 = [
    ("Do they make sense?", PEOPLE, [
        "Which principles resonate most for our community?",
        "Where do they fall short, or feel abstract?",
        "Whose voice is centred — and whose is missing?",
    ]),
]
rect(s, qx, 2.25, cw, 4.0, fill=NAVY2, rounded=True, radius=0.04)
tb(s, qx + 0.4, 2.55, cw - 0.8, 0.5,
   [{'t': "DO THEY MAKE SENSE?", 'size': 14, 'bold': True, 'color': PEOPLE, 'font': FONT_H}])
yy = 3.15
for it in ["Which principles resonate most for our community?",
           "Where do they fall short, or feel abstract?",
           "Whose voice is centred — and whose is missing?",
           "Does 'equalising power' go far enough for families?"]:
    tb(s, qx + 0.4, yy, cw - 0.8, 0.7,
       [{'t': [("?  ", {'color': PEOPLE, 'bold': True}), (it, {'color': RGBColor(0xD7,0xDE,0xE9)})],
         'size': 14.5, 'font': FONT_B, 'line': 1.18}])
    yy += 0.74
# right: prompts for what's missing
qx2 = qx + cw + 0.5
cw2 = SW - MX - qx2
rect(s, qx2, 2.25, cw2, 4.0, fill=PEOPLE, rounded=True, radius=0.04)
tb(s, qx2 + 0.4, 2.55, cw2 - 0.8, 0.5,
   [{'t': "WHAT MIGHT BE MISSING FOR RARE DISEASE?", 'size': 14, 'bold': True,
     'color': INK, 'font': FONT_H}])
missing = [
    ("Urgency", "progressive & life-limiting conditions — time matters"),
    ("Diagnostic odyssey", "designing amid uncertainty and 'no diagnosis yet'"),
    ("Ultra-rare / n-of-few", "when the 'community' is a handful of families"),
    ("Distance & equity", "rural & regional NSW; tele-everything"),
    ("Trust & data", "research fatigue; who owns our stories and data"),
    ("Sustainability", "small, volunteer-run groups; carer burnout"),
]
yy = 3.12
for t, d in missing:
    tb(s, qx2 + 0.4, yy, cw2 - 0.8, 0.55,
       [{'t': [("+  ", {'bold': True, 'color': PRACTICE}),
               (t + "  —  ", {'bold': True, 'color': INK}),
               (d, {'color': RGBColor(0x4A,0x3D,0x12)})],
         'size': 13, 'font': FONT_B, 'line': 1.12}])
    yy += 0.5
page_no(s, num())

# ============================================================================
# SLIDE 10 — CODE OF CARE (rules of engagement)
# ============================================================================
s = slide()
kicker_title(s, "Rules of engagement", "Our Code of Care", accent=PRACTICE)
tb(s, MX, 1.95, SW - 2 * MX, 0.6,
   [{'t': "How we'll be together today — and beyond. We'll add to and adopt these as a group.",
     'size': 15.5, 'color': INK2, 'font': FONT_B, 'line': 1.25}])
care = [
    ("Nothing about us, without us", "Lived experience leads. People most affected help shape every decision."),
    ("Every voice counts", "All knowledge is valued. No jargon left unexplained; ask anything."),
    ("Share the air", "Step up, step back. Make room so quieter voices are heard."),
    ("Safe & confidential", "Stories stay in the room. Care for yourself and each other."),
    ("Disagree generously", "Challenge ideas, not people. Curiosity over certainty."),
    ("It's okay to pause", "Distress is welcome here; support is available. Progress over perfection."),
]
cw = (SW - 2 * MX - 2 * 0.3) / 3
ch = 1.7
positions = []
for r in range(2):
    for col in range(3):
        positions.append((MX + col * (cw + 0.3), 2.7 + r * (ch + 0.28)))
for i, ((t, d), (x, y)) in enumerate(zip(care, positions)):
    rect(s, x, y, cw, ch, fill=PRACTICE_T, rounded=True, radius=0.06)
    oval(s, x + 0.3, y + 0.28, 0.46, 0.46, fill=PRACTICE)
    tb(s, x + 0.3, y + 0.28, 0.46, 0.46,
       [{'t': str(i + 1), 'size': 16, 'bold': True, 'color': WHITE,
         'align': PP_ALIGN.CENTER, 'font': FONT_H}], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, x + 0.92, y + 0.26, cw - 1.1, 0.6,
       [{'t': t, 'size': 14.5, 'bold': True, 'color': INK, 'font': FONT_H, 'line': 1.0}])
    tb(s, x + 0.3, y + 0.92, cw - 0.55, 0.7,
       [{'t': d, 'size': 11.8, 'color': INK2, 'font': FONT_B, 'line': 1.18}])
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 11 — THE FOUR DOMAINS (transition)
# ============================================================================
s = slide(PAPER_ALT)
tb(s, MX, 0.85, 11.0, 0.4,
   [{'t': "WORKING THROUGH THE FRAMEWORK", 'size': 13.5, 'bold': True, 'color': NAVY, 'font': FONT_H}])
tb(s, MX, 1.28, 11.6, 0.7,
   [{'t': "Four domains to design with", 'size': 31, 'bold': True, 'color': INK, 'font': FONT_H}])
dom = [
    ("People", "Who needs to be involved?", PEOPLE, INK, ["Self", "Team", "Organisation", "Community"]),
    ("Process", "How will we approach this?", PROCESS, WHITE, ["Invitation", "Preparation", "Navigation", "Completion"]),
    ("Practice", "What will we use?", PRACTICE, WHITE, ["Tools", "Techniques", "Methods", "Principles"]),
    ("Place", "Where does this fit?", PLACE, WHITE, ["Symbols", "Products", "Interactions", "System"]),
]
cw = (SW - 2 * MX - 3 * 0.3) / 4
x = MX
for name, q, c, tc, items in dom:
    top = 2.5
    rect(s, x, top, cw, 3.7, fill=c, rounded=True, radius=0.05)
    tb(s, x + 0.3, top + 0.32, cw - 0.6, 0.6,
       [{'t': name, 'size': 23, 'bold': True, 'color': tc, 'font': FONT_H}])
    tb(s, x + 0.3, top + 0.95, cw - 0.6, 0.7,
       [{'t': q, 'size': 13.5, 'italic': True,
         'color': tc if tc == WHITE else RGBColor(0x5A,0x47,0x12), 'font': FONT_B, 'line': 1.15}])
    line_h(s, x + 0.3, top + 1.75, cw - 0.6,
           color=(RGBColor(0xFF,0xFF,0xFF) if tc == WHITE else INK), weight=0.75)
    yy = top + 1.95
    for it in items:
        tb(s, x + 0.3, yy, cw - 0.6, 0.4,
           [{'t': it, 'size': 13.5, 'bold': True, 'color': tc, 'font': FONT_B}])
        yy += 0.42
    x += cw + 0.3
tb(s, MX, 6.4, 11.6, 0.4,
   [{'t': "We'll look closely at People — the heart of co-design — then Process, Practice and Place.",
     'size': 13, 'italic': True, 'color': INK2, 'font': FONT_B}])
page_no(s, num())

# ============================================================================
# SLIDE 12 — PEOPLE: who needs to be involved
# ============================================================================
s = slide()
kicker_title(s, "People · domain 1", "Who needs to be involved?", accent=PEOPLE, title_color=INK)
tb(s, MX, 1.95, 5.0, 1.4,
   [{'t': "Co-design ripples outward — from how we show up ourselves, to the team, "
          "the organisation, and the wider community we serve.",
     'size': 15.5, 'color': INK2, 'font': FONT_B, 'line': 1.3}])
tb(s, MX, 3.5, 5.0, 1.2,
   [{'t': [("For rare disease: ", {'bold': True, 'color': PRACTICE}),
           ("the 'community' may be geographically scattered and tiny — and patients, "
            "carers and families are experts in equal measure.", {'color': INK2})],
     'size': 14, 'font': FONT_B, 'line': 1.28}])
# concentric rings (right)
cx, cy = 9.6, 3.95
rings = [
    (3.4, PEOPLE_T, "Community", INK, 0.0),
    (2.55, RGBColor(0xF2,0xDF,0xA6), "Organisation", INK, 0.62),
    (1.7, RGBColor(0xEC,0xC9,0x55), "Team", INK, 0.62),
    (0.95, PEOPLE, "Self", INK, 0.0),
]
for d, col, lbl, tc, lblpos in rings:
    oval(s, cx - d / 2, cy - d / 2, d, d, fill=col, line=WHITE, line_w=2.0)
for d, col, lbl, tc, lblpos in rings:
    if lbl == "Self":
        tb(s, cx - 0.5, cy - 0.18, 1.0, 0.4,
           [{'t': lbl, 'size': 13, 'bold': True, 'color': tc, 'align': PP_ALIGN.CENTER, 'font': FONT_H}])
    else:
        tb(s, cx - d / 2, cy - d / 2 + 0.12, d, 0.4,
           [{'t': lbl, 'size': 13, 'bold': True, 'color': tc, 'align': PP_ALIGN.CENTER, 'font': FONT_H}])
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 13 — PEOPLE: key teams
# ============================================================================
s = slide()
kicker_title(s, "People · key teams", "The teams that make it work", accent=PEOPLE, title_color=INK)
# nested circles (left)
cx, cy = 4.0, 4.25
oval(s, cx - 1.85, cy - 1.85, 3.7, 3.7, fill=EXPLORE)
oval(s, cx - 1.3, cy - 1.3 + 0.25, 2.6, 2.6, fill=PRACTICE)
oval(s, cx - 0.75, cy - 0.75 + 0.45, 1.5, 1.5, fill=PLACE)
# co-design team = translucent "lens" that overlaps the inner teams, bottom-right
oval(s, cx + 0.05, cy + 0.2, 1.9, 1.9, fill=WHITE, alpha=0.62, line=PRACTICE, line_w=2.5)
tb(s, cx - 1.8, cy - 1.7, 3.6, 0.4, [{'t': "Community", 'size': 13, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER, 'font': FONT_H}])
tb(s, cx - 1.3, cy - 1.0, 2.6, 0.4, [{'t': "Organisation", 'size': 12, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER, 'font': FONT_H}])
tb(s, cx - 0.7, cy + 0.18, 1.4, 0.4, [{'t': "Project team", 'size': 11, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER, 'font': FONT_H}], anchor=MSO_ANCHOR.MIDDLE)
tb(s, cx + 0.62, cy + 1.18, 1.7, 0.4, [{'t': "Co-design team", 'size': 11, 'bold': True, 'color': INK, 'align': PP_ALIGN.CENTER, 'font': FONT_H}], anchor=MSO_ANCHOR.MIDDLE)
# explanation (right)
ex = 6.7
items = [
    ("Community", "Everyone affected by, and with a stake in, the issue.", EXPLORE),
    ("Organisation", "Those who hold resources, mandate and accountability.", PRACTICE),
    ("Project team", "The core group driving the work day to day.", PLACE),
    ("Co-design team", "Where lived experience and the project team meet — the engine of co-design.", PRACTICE),
]
yy = 2.15
for t, d, c in items:
    oval(s, ex, yy + 0.06, 0.22, 0.22, fill=c)
    tb(s, ex + 0.4, yy - 0.04, SW - MX - ex - 0.4, 0.5,
       [{'t': t, 'size': 16, 'bold': True, 'color': INK, 'font': FONT_H, 'sa': 1},
        {'t': d, 'size': 13, 'color': INK2, 'font': FONT_B, 'line': 1.2}])
    yy += 1.05
source_note(s, "After Blomkamp / New Know How: Co-Design Teams.")
page_no(s, num())

# ============================================================================
# SLIDE 14 — PEOPLE: key roles
# ============================================================================
s = slide()
kicker_title(s, "People · key roles", "The roles around the table", accent=PEOPLE, title_color=INK)
roles = [
    ("Convenor", "facilitator", PEOPLE),
    ("Creator", "designer", PROCESS),
    ("Commissioner", "sponsor", PRACTICE),
    ("Connector", "broker", PLACE),
    ("Coach", "mentor", PEOPLE),
    ("Living expert", "lived experience", PROCESS),
    ("Professional expert", "practice wisdom", PRACTICE),
    ("Carer", "support worker", PLACE),
    ("Inquirer", "critical friend", PEOPLE),
    ("Coordinator", "producer", PROCESS),
]
cw = (SW - 2 * MX - 4 * 0.22) / 5
ch = 1.55
for i, (t, d, c) in enumerate(roles):
    col = i % 5
    row = i // 5
    x = MX + col * (cw + 0.22)
    y = 2.2 + row * (ch + 0.28)
    rect(s, x, y, cw, ch, fill=PAPER_ALT, rounded=True, radius=0.08)
    rect(s, x, y, cw, 0.11, fill=c)
    tb(s, x + 0.18, y + 0.34, cw - 0.36, 0.7,
       [{'t': t, 'size': 14.5, 'bold': True, 'color': INK, 'font': FONT_H, 'line': 1.0}])
    tb(s, x + 0.18, y + 1.02, cw - 0.36, 0.4,
       [{'t': d, 'size': 11.5, 'italic': True, 'color': INK2, 'font': FONT_B}])
tb(s, MX, 6.05, SW - 2 * MX, 0.5,
   [{'t': [("One person may hold several roles — and ", {'color': INK2}),
           ("living experts, carers and professional experts sit as equals.",
            {'bold': True, 'color': INK})],
     'size': 13.5, 'italic': True, 'font': FONT_B}])
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 15 — PEOPLE: shades of co-design (matrix)
# ============================================================================
s = slide()
kicker_title(s, "People · depth of participation", "Shades of co-design", accent=PEOPLE, title_color=INK)
tb(s, MX, 1.9, SW - 2 * MX, 0.5,
   [{'t': "How deeply do people participate — and across which kinds of knowledge? "
          "Different decisions sit at different shades.",
     'size': 14.5, 'color': INK2, 'font': FONT_B, 'line': 1.2}])
depth = [
    ("Talk to", RGBColor(0xD7,0xF0,0xDE)),
    ("Listen & respond to", RGBColor(0xAE,0xE2,0xBE)),
    ("Collaborate with", RGBColor(0x74,0xCE,0x92)),
    ("Make decisions together", RGBColor(0x3F,0xB9,0x6B)),
    ("Delegate authority to", RGBColor(0x1F,0xA6,0x4A)),
]
knowledge = ["Lived\nexperience", "Practice\nwisdom", "Research\nevidence", "Cultural\nknowledge"]
gx, gy = MX, 2.6
labelw = 3.0
colw = (SW - 2 * MX - labelw) / 4
rowh = 0.66
# column headers
for j, k in enumerate(knowledge):
    tb(s, gx + labelw + j * colw, gy - 0.02, colw, 0.5,
       [{'t': k.replace("\n", " "), 'size': 11.5, 'bold': True, 'color': INK,
         'align': PP_ALIGN.CENTER, 'font': FONT_H, 'line': 0.95}])
gy0 = gy + 0.5
for i, (dlabel, dcol) in enumerate(depth):
    y = gy0 + i * rowh
    rect(s, gx, y, labelw - 0.12, rowh - 0.1, fill=dcol, rounded=True, radius=0.12)
    tb(s, gx + 0.2, y, labelw - 0.5, rowh - 0.1,
       [{'t': dlabel, 'size': 12.5, 'bold': True,
         'color': (WHITE if i >= 3 else INK), 'font': FONT_H}], anchor=MSO_ANCHOR.MIDDLE)
    for j in range(4):
        cxx = gx + labelw + j * colw
        rect(s, cxx + 0.06, y, colw - 0.12, rowh - 0.1, fill=PAPER_ALT, line=LINE,
             line_w=1.0, rounded=True, radius=0.12)
        oval(s, cxx + colw / 2 - 0.11, y + (rowh - 0.1) / 2 - 0.11, 0.22, 0.22,
             fill=PAPER, line=dcol, line_w=2.0)
tb(s, gx, gy0 + 5 * rowh + 0.08, SW - 2 * MX, 0.4,
   [{'t': [("Deeper isn't always better — ", {'bold': True, 'color': INK}),
           ("match the shade to the decision, the stakes and what people have capacity for.",
            {'color': INK2})], 'size': 12.5, 'italic': True, 'font': FONT_B}])
source_note(s, "After Blomkamp (2025), Shades of Co-Design.")
page_no(s, num())

# ============================================================================
# SLIDE 16 — PROCESS
# ============================================================================
s = slide()
kicker_title(s, "Process · domain 2", "How will we approach this?", accent=PROCESS)
tb(s, MX, 1.95, SW - 2 * MX, 0.5,
   [{'t': "Co-design moves through phases — each needs care, not just the 'doing' in the middle.",
     'size': 15.5, 'color': INK2, 'font': FONT_B, 'line': 1.2}])
phases = [
    ("Invitation", "Who we ask, how, and why — the first act of power-sharing.",
     "Reach scattered families early; make it genuinely easy to say yes."),
    ("Preparation", "Building trust, safety and readiness before the work begins.",
     "Plain-language info; flexible formats; support to take part."),
    ("Navigation", "Doing the work together, adapting as we learn.",
     "Hold space for distress, complexity and changing health."),
    ("Completion", "Closing well, sharing back, and sustaining relationships.",
     "Show what changed; don't disappear — keep the loop open."),
]
cw = (SW - 2 * MX - 3 * 0.28) / 4
x = MX
for i, (t, d, r) in enumerate(phases):
    top = 2.65
    rect(s, x, top, cw, 3.5, fill=PROCESS_T, rounded=True, radius=0.05)
    oval(s, x + 0.28, top + 0.3, 0.5, 0.5, fill=PROCESS)
    tb(s, x + 0.28, top + 0.3, 0.5, 0.5,
       [{'t': str(i + 1), 'size': 17, 'bold': True, 'color': WHITE,
         'align': PP_ALIGN.CENTER, 'font': FONT_H}], anchor=MSO_ANCHOR.MIDDLE)
    tb(s, x + 0.28, top + 1.0, cw - 0.5, 0.5,
       [{'t': t, 'size': 17, 'bold': True, 'color': INK, 'font': FONT_H}])
    tb(s, x + 0.28, top + 1.5, cw - 0.56, 1.0,
       [{'t': d, 'size': 12.5, 'color': INK, 'font': FONT_B, 'line': 1.22}])
    line_h(s, x + 0.28, top + 2.55, cw - 0.56, color=PROCESS, weight=0.75)
    tb(s, x + 0.28, top + 2.68, cw - 0.56, 0.8,
       [{'t': [("Rare disease: ", {'bold': True, 'color': PROCESS}), (r, {'color': INK2})],
         'size': 11, 'font': FONT_B, 'line': 1.18}])
    x += cw + 0.28
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 17 — PRACTICE
# ============================================================================
s = slide()
kicker_title(s, "Practice · domain 3", "What will we use?", accent=PRACTICE)
tb(s, MX, 1.95, 5.2, 1.6,
   [{'t': "Our toolkit, from the concrete to the conceptual. Choose methods that fit "
          "the people and the moment — not the other way around.",
     'size': 15.5, 'color': INK2, 'font': FONT_B, 'line': 1.3}])
tb(s, MX, 3.7, 5.2, 1.6,
   [{'t': [("For rare disease, favour: ", {'bold': True, 'color': PRACTICE}),
           ("low-bandwidth, remote-friendly, asynchronous methods; storytelling; and "
            "tools families can use from home, in their own time.", {'color': INK2})],
     'size': 14, 'font': FONT_B, 'line': 1.28}])
layers = [
    ("Tools", "The tangible things we pick up — templates, canvases, prompts, prototypes.", PRACTICE),
    ("Techniques", "Specific moves — interviewing, journey mapping, sense-making.", RGBColor(0xEE,0x77,0x40)),
    ("Methods", "Bundled approaches — participatory design, action research.", RGBColor(0xF2,0x96,0x6A)),
    ("Principles", "The why beneath the how — what good practice looks like.", RGBColor(0xF7,0xB5,0x97)),
]
bx = 6.5
bw0 = SW - MX - bx
yy = 2.0
for i, (t, d, c) in enumerate(layers):
    h = 1.05
    inset = i * 0.0
    rect(s, bx, yy, bw0, h, fill=c, rounded=True, radius=0.05)
    tb(s, bx + 0.32, yy + 0.16, bw0 - 0.6, 0.4,
       [{'t': t, 'size': 16.5, 'bold': True, 'color': (WHITE if i < 2 else INK), 'font': FONT_H}])
    tb(s, bx + 0.32, yy + 0.56, bw0 - 0.6, 0.5,
       [{'t': d, 'size': 11.8, 'color': (WHITE if i < 2 else INK), 'font': FONT_B, 'line': 1.1}])
    yy += h + 0.12
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 18 — PLACE
# ============================================================================
s = slide()
kicker_title(s, "Place · domain 4", "Where does this fit?", accent=PLACE)
tb(s, MX, 1.95, 5.2, 1.8,
   [{'t': "Co-design doesn't happen in a vacuum. 'Place' asks where our work lands — "
          "from the smallest signal to the whole system — so change actually sticks.",
     'size': 15.5, 'color': INK2, 'font': FONT_B, 'line': 1.3}])
tb(s, MX, 4.0, 5.2, 1.8,
   [{'t': [("For rare disease, the 'system' is fragmented — ", {'bold': True, 'color': PLACE}),
           ("health, genetics, disability, NDIS, education, research and community groups. "
            "Place asks how our work connects them.", {'color': INK2})],
     'size': 14, 'font': FONT_B, 'line': 1.28}])
levels = [
    ("System", "Policies, funding, structures and the way the whole field works.", PLACE),
    ("Interactions", "How people and services meet — the moments of contact and care.", RGBColor(0x91,0x66,0xE2)),
    ("Products", "The services, tools and outputs people actually use.", RGBColor(0xAC,0x8B,0xEC)),
    ("Symbols", "Language, stories and meaning — what things signal and stand for.", RGBColor(0xC9,0xB4,0xF3)),
]
bx = 6.5
yy = 2.0
fullw = SW - MX - bx
for i, (t, d, c) in enumerate(levels):
    w = fullw - i * 0.7
    x = bx + i * 0.35
    h = 1.02
    rect(s, x, yy, w, h, fill=c, rounded=True, radius=0.06)
    tb(s, x + 0.3, yy + 0.14, w - 0.5, 0.4,
       [{'t': t, 'size': 16, 'bold': True, 'color': (WHITE if i < 2 else INK), 'font': FONT_H}])
    tb(s, x + 0.3, yy + 0.54, w - 0.5, 0.5,
       [{'t': d, 'size': 11.3, 'color': (WHITE if i < 2 else INK), 'font': FONT_B, 'line': 1.1}])
    yy += h + 0.1
footer_brand(s); page_no(s, num())

# ============================================================================
# SLIDE 19 — CO-DESIGN MATURITY MODEL
# ============================================================================
s = slide()
kicker_title(s, "Where are we?", "The co-design maturity journey", accent=NAVY)
mat = [
    ("Flourish", "Evolving and spreading co-design across the system.", FLOURISH, INK, 5.4),
    ("Integrate", "Adapting and embedding co-design in how we work.", INTEGRATE, WHITE, 6.0),
    ("Apply", "Using co-design principles and practices in context.", APPLY, WHITE, 6.6),
    ("Understand", "Identifying the main parts of co-design and the value it delivers.", UNDERSTAND, WHITE, 7.2),
    ("Explore", "Getting curious about co-design.", EXPLORE, WHITE, 7.8),
]
top = 2.15
cx = MX + 3.9
rowh = 0.82
for i, (t, d, c, tc, w) in enumerate(mat):
    y = top + i * rowh
    x = cx - w / 2
    rect(s, x, y, w, rowh - 0.12, fill=c, rounded=False)
    tb(s, x, y, w, rowh - 0.12,
       [{'t': t, 'size': 16, 'bold': True, 'color': tc, 'align': PP_ALIGN.CENTER, 'font': FONT_H}],
       anchor=MSO_ANCHOR.MIDDLE)
    # description to the right
    tb(s, cx + 4.1, y + 0.04, SW - MX - (cx + 4.1), rowh - 0.12,
       [{'t': d, 'size': 12.5, 'color': INK2, 'font': FONT_B, 'line': 1.1}],
       anchor=MSO_ANCHOR.MIDDLE)
tb(s, MX, top + 5 * rowh + 0.05, 7.2, 0.5,
   [{'t': [("Quick self-check: ", {'bold': True, 'color': INK}),
           ("where is Rare Diseases NSW today — and where do we want to be?",
            {'italic': True, 'color': INK2})], 'size': 13.5, 'font': FONT_B}])
source_note(s, "After Blomkamp / New Know How: Co-Design Maturity Model.")
page_no(s, num())

# ============================================================================
# SLIDE 20 — HOW TO / NEXT STEPS
# ============================================================================
s = slide(NAVY)
rect(s, 0, 0, SW, 0.16, fill=PROCESS)
tb(s, MX, 0.62, 11.0, 0.4,
   [{'t': "HOW TO — TAKING THIS FORWARD", 'size': 13.5, 'bold': True, 'color': PEOPLE, 'font': FONT_H}])
tb(s, MX, 1.02, 11.8, 0.9,
   [{'t': "First steps for Rare Diseases NSW", 'size': 30, 'bold': True, 'color': WHITE,
     'font': FONT_H}])
steps = [
    ("Capture today", "Turn our input into a 'Rare Diseases NSW co-design framework v0.1'."),
    ("Form the team", "Confirm the co-design team and who holds which role."),
    ("Choose a first project", "Agree one real, bounded challenge to pilot together."),
    ("Set our shades", "Decide who's consulted, who collaborates, who decides."),
    ("Resource it", "Fund and pay lived-experience contributors; build in support."),
    ("Learn & adapt", "Schedule check-ins; revisit the framework as we go."),
]
cw = (SW - 2 * MX - 2 * 0.3) / 3
ch = 1.7
for i, (t, d) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = MX + col * (cw + 0.3)
    y = 2.35 + row * (ch + 0.3)
    rect(s, x, y, cw, ch, fill=NAVY2, rounded=True, radius=0.06)
    tb(s, x + 0.32, y + 0.26, cw - 0.6, 0.5,
       [{'t': f"0{i+1}", 'size': 15, 'bold': True, 'color': PEOPLE, 'font': FONT_H}])
    tb(s, x + 0.32, y + 0.66, cw - 0.6, 0.5,
       [{'t': t, 'size': 16, 'bold': True, 'color': WHITE, 'font': FONT_H}])
    tb(s, x + 0.32, y + 1.1, cw - 0.6, 0.6,
       [{'t': d, 'size': 12, 'color': RGBColor(0xC2,0xCC,0xDA), 'font': FONT_B, 'line': 1.18}])
page_no(s, num())

# ============================================================================
# SLIDE 21 — THANK YOU / CLOSE
# ============================================================================
s = slide(NAVY)
bx, bw = 0, SW / 4
for c in (PEOPLE, PROCESS, PRACTICE, PLACE):
    rect(s, bx, SH - 0.16, bw, 0.16, fill=c); bx += bw
tb(s, MX, 2.1, 11.6, 1.2,
   [{'t': "Thank you", 'size': 54, 'bold': True, 'color': WHITE, 'font': FONT_H}])
tb(s, MX, 3.35, 11.0, 1.0,
   [{'t': [("Designed ", {'color': RGBColor(0xCE,0xD6,0xE2)}),
           ("with", {'bold': True, 'color': PEOPLE}),
           (" the rare disease community — not for it.", {'color': RGBColor(0xCE,0xD6,0xE2)})],
     'size': 20, 'font': FONT_B, 'line': 1.2}])
line_h(s, MX, 4.55, 4.2, color=NAVY2, weight=1.5)
tb(s, MX, 4.8, 11.0, 1.2,
   [{'t': [("Rare NSW · Kris Pierce Consulting", {'bold': True, 'color': WHITE})],
     'size': 15, 'font': FONT_B, 'sa': 3},
    {'t': "Framework adapted from Emma Blomkamp's Systemic Design Practice & New Know How "
          "(Shades of Co-Design, Co-Design Teams, Maturity Model). CC BY-NC.",
     'size': 11.5, 'color': RGBColor(0x9D,0xAA,0xBD), 'font': FONT_B, 'line': 1.25}])

# ----------------------------------------------------------------------------
# Output next to this script so the generator is portable/reproducible.
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Rare-Diseases-NSW-Co-Design-Workshop.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
